# -*- coding: utf-8 -*-
"""
render_pptx.py — рендер тизера в стиле Big 3 (McKinsey/Strategy Partners).
Реализует дизайн-систему из presentation_design.md: kicker + action-title +
акцентные подзаголовки + чистые чарты + футер (источник · конфиденциально · №).
Белый фон, сдержанная палитра, БЕЗ фото в теле. Тема в THEME -> перекраска под бренд.

Зависимости: python-pptx, Pillow (для обрезки логотипа).
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

THEME = {  # McKinsey-инспирированная палитра; меняется под бренд клиента
    "primary": RGBColor(0x04, 0x24, 0x4A),   # глубокий navy
    "accent":  RGBColor(0x00, 0xA9, 0xCE),   # cyan/teal
    "accent2": RGBColor(0xE8, 0xA3, 0x3D),   # тёплый (только для контраста)
    "ink":     RGBColor(0x1A, 0x1A, 0x1A),
    "muted":   RGBColor(0x8A, 0x94, 0xA6),
    "hairline": RGBColor(0xE3, 0xE8, 0xEE),
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "font":    "Arial",
}
W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)


def _pres():
    p = Presentation(); p.slide_width = W; p.slide_height = H
    return p


def _blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def _rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _txt(slide, l, t, w, h, text, size=11, color=None, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, caps=False, spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if spacing:
            para.space_after = Pt(spacing)
        bullet = isinstance(ln, tuple)
        s = (ln[0] if bullet else ln) or ""
        if caps:
            s = s.upper()
        run = para.add_run(); run.text = ("—  " + s) if bullet else s
        run.font.size = Pt(size); run.font.bold = bold; run.font.name = THEME["font"]
        run.font.color.rgb = color or THEME["ink"]
    return tb


def _kicker_title(slide, kicker, action_title):
    _txt(slide, MARGIN, Inches(0.3), W - 2*MARGIN, Inches(0.3),
         kicker, size=10, color=THEME["muted"], bold=True, caps=True)
    _txt(slide, MARGIN, Inches(0.62), W - 2*MARGIN, Inches(0.95),
         action_title, size=20, color=THEME["primary"], bold=True)
    _rect(slide, MARGIN, Inches(1.62), W - 2*MARGIN, Pt(1.5), THEME["hairline"])


def _subhead(slide, l, t, w, text):
    """Акцентная плашка-подзаголовок над блоком данных (как в Strategy Partners)."""
    _rect(slide, l, t, w, Inches(0.34), THEME["accent"])
    _txt(slide, l + Inches(0.12), t + Inches(0.04), w - Inches(0.2), Inches(0.28),
         text, size=12, color=THEME["white"], bold=True)


def _footer(slide, source="", page=None):
    _rect(slide, 0, Inches(7.12), W, Pt(1), THEME["hairline"])
    if source:
        _txt(slide, MARGIN, Inches(7.18), Inches(8), Inches(0.25),
             "Источник: " + source, size=8, color=THEME["muted"])
    _txt(slide, Inches(5.5), Inches(7.18), Inches(2.3), Inches(0.25),
         "Конфиденциально", size=8, color=THEME["muted"], align=PP_ALIGN.CENTER)
    if page:
        _txt(slide, W - MARGIN - Inches(1), Inches(7.18), Inches(1), Inches(0.25),
             str(page), size=8, color=THEME["muted"], align=PP_ALIGN.RIGHT)


def _pic_fit(slide, path, l, t, w, h):
    """Картинка crop-to-fit (логотип/актив). Используется РЕДКО (обложка)."""
    if not (path and Path(str(path)).exists()):
        return False
    try:
        from PIL import Image
        iw, ih = Image.open(str(path)).size
        tr, ir = w / h, iw / ih
        pic = slide.shapes.add_picture(str(path), l, t, width=w, height=h)
        if ir > tr:
            cut = (1 - tr / ir) / 2; pic.crop_left = cut; pic.crop_right = cut
        else:
            cut = (1 - ir / tr) / 2; pic.crop_top = cut; pic.crop_bottom = cut
        return True
    except Exception:
        try:
            slide.shapes.add_picture(str(path), l, t, width=w, height=h); return True
        except Exception:
            return False


def _chart(slide, path, l, t, w, h):
    if path and Path(str(path)).exists():
        try:
            from PIL import Image
            iw, ih = Image.open(str(path)).size
            scale = min(w / iw, h / ih)
            cw, ch = int(iw * scale), int(ih * scale)
            slide.shapes.add_picture(str(path), l + (w - cw)//2, t, width=cw, height=ch)
            return True
        except Exception:
            slide.shapes.add_picture(str(path), l, t, width=w); return True
    return False


# ===================== СЛАЙДЫ =====================
def slide_cover(p, c):
    s = _blank(p)
    _rect(s, 0, 0, W, H, THEME["primary"])
    hero = c["cover"].get("hero")
    if hero and Path(str(hero)).exists():
        # full-bleed герой-фото (релевантное, проверенное зрением) + тёмная плашка под текст
        _pic_fit(s, hero, 0, 0, W, H)
        _rect(s, 0, Inches(4.6), W, H - Inches(4.6), THEME["primary"], line=None)
        _overlay_alpha(s.shapes[-1], 78)  # полупрозрачная navy-плашка
    logo = c["cover"].get("logo")
    if logo and Path(str(logo)).exists():
        _pic_fit(s, logo, MARGIN, Inches(0.55), Inches(2.4), Inches(0.85))
    _rect(s, MARGIN, Inches(4.9), Inches(1.4), Pt(3), THEME["accent"])
    _txt(s, MARGIN, Inches(5.15), Inches(11.5), Inches(1.4),
         c["cover"]["positioning"], size=26, color=THEME["white"], bold=True)
    _txt(s, MARGIN, Inches(6.55), Inches(8), Inches(0.4),
         c["cover"]["subtitle"], size=14, color=THEME["accent"], bold=True, caps=True)
    _txt(s, W - MARGIN - Inches(3), Inches(6.55), Inches(3), Inches(0.4),
         c["cover"].get("date", ""), size=11, color=THEME["white"], align=PP_ALIGN.RIGHT)
    return s


def _overlay_alpha(shape, percent):
    """Полупрозрачная заливка фигуры (python-pptx не имеет API -> правим XML)."""
    try:
        from pptx.oxml.ns import qn
        sp = shape.fill._xPr.find(qn('a:solidFill'))
        srgb = sp.find(qn('a:srgbClr'))
        alpha = sp.makeelement(qn('a:alpha'), {'val': str(int(percent * 1000))})
        srgb.append(alpha)
    except Exception:
        pass


def slide_exec_summary(p, c):
    s = _blank(p)
    sm = c["summary"]
    _kicker_title(s, c["company_name"], sm["action_title"])
    # левая колонка: тезис + хайлайты
    _txt(s, MARGIN, Inches(1.85), Inches(7.4), Inches(0.9), sm["one_liner"],
         size=12, color=THEME["ink"], spacing=4)
    _subhead(s, MARGIN, Inches(2.95), Inches(7.4), "Почему это привлекательно")
    _txt(s, MARGIN, Inches(3.45), Inches(7.4), Inches(3.4),
         [(h,) for h in sm["highlights"][:6]], size=11.5, color=THEME["ink"], spacing=6)
    # правая карточка: ключевые цифры
    cx = Inches(8.6); cw = Inches(4.1)
    _rect(s, cx, Inches(1.85), cw, Inches(5.0), THEME["primary"])
    _txt(s, cx + Inches(0.25), Inches(2.05), cw - Inches(0.5), Inches(0.4),
         "Ключевые показатели", size=13, color=THEME["accent"], bold=True, caps=True)
    kf = sm["key_figures"]; y = 2.65
    items = [("Выручка (2025)", kf.get("revenue")), ("EBITDA (норм.)", kf.get("ebitda")),
             ("Клиентская база", kf.get("clients")), ("Оценка EV", kf.get("ev_range"))]
    for label, val in items:
        _txt(s, cx + Inches(0.25), Inches(y), cw - Inches(0.5), Inches(0.3),
             label, size=10, color=THEME["muted"], caps=True)
        _txt(s, cx + Inches(0.25), Inches(y + 0.28), cw - Inches(0.5), Inches(0.5),
             str(val), size=18, color=THEME["white"], bold=True)
        y += 1.05
    _footer(s, sm.get("source", "ФНС / ГИР БО; анализ"), 2)
    return s


def slide_company(p, c):
    s = _blank(p); co = c["company"]
    _kicker_title(s, c["company_name"], co["action_title"])
    _subhead(s, MARGIN, Inches(1.85), Inches(5.9), "Профиль компании")
    rows = [("Основана", co.get("founded")), ("География", co.get("location")),
            ("Команда", co.get("team_size")), ("Бренды", ", ".join(co.get("brands", []))),
            ("Модель", co.get("business_model"))]
    y = 2.45
    for k, v in rows:
        _txt(s, MARGIN, Inches(y), Inches(1.7), Inches(0.5), k, size=10.5,
             color=THEME["muted"], bold=True, caps=True)
        _txt(s, MARGIN + Inches(1.8), Inches(y), Inches(4.0), Inches(0.7), str(v),
             size=11.5, color=THEME["ink"])
        y += 0.78
    _subhead(s, Inches(6.9), Inches(1.85), Inches(5.8), "Эксклюзивы и периметр сделки")
    _txt(s, Inches(6.9), Inches(2.45), Inches(5.8), Inches(3.5),
         [(x,) for x in co.get("exclusives", [])] +
         [(("Периметр: " + ", ".join(co.get("deal_perimeter", []))),)],
         size=11.5, color=THEME["ink"], spacing=6)
    _footer(s, "ЕГРЮЛ / ГИР БО; данные компании", 3)
    return s


def slide_market(p, c, chart=None):
    s = _blank(p); mk = c["market"]
    _kicker_title(s, c["company_name"], mk["action_title"])
    _subhead(s, MARGIN, Inches(1.85), Inches(6.0), "Динамика рынка")
    _chart(s, chart, MARGIN, Inches(2.35), Inches(6.0), Inches(3.4))
    _subhead(s, Inches(7.0), Inches(1.85), Inches(5.7), "Драйверы и риски")
    _txt(s, Inches(7.0), Inches(2.45), Inches(5.7), Inches(2.8),
         [(x,) for x in mk.get("tailwinds", [])], size=11.5, color=THEME["ink"], spacing=6)
    if mk.get("risks"):
        _txt(s, Inches(7.0), Inches(5.2), Inches(5.7), Inches(1.3),
             [("Риски: " + "; ".join(mk["risks"]),)], size=10, color=THEME["muted"])
    _footer(s, "; ".join(mk.get("sources", ["отраслевые отчёты"])[:3]), 4)
    return s


def slide_financials(p, c, chart_rev=None, chart_margin=None):
    s = _blank(p)
    _kicker_title(s, c["company_name"], c["financials"]["action_title"])
    _subhead(s, MARGIN, Inches(1.85), Inches(5.9), "Выручка")
    _chart(s, chart_rev, MARGIN, Inches(2.35), Inches(5.9), Inches(2.7))
    _subhead(s, Inches(6.9), Inches(1.85), Inches(5.8), "Прибыль и маржинальность")
    _chart(s, chart_margin, Inches(6.9), Inches(2.35), Inches(5.8), Inches(2.7))
    _subhead(s, MARGIN, Inches(5.3), Inches(12.1), "Комментарий")
    _txt(s, MARGIN, Inches(5.8), Inches(12.1), Inches(1.2),
         c["financials"].get("commentary", ""), size=11, color=THEME["ink"], spacing=3)
    _footer(s, "ГИР БО (РСБУ); анализ", 5)
    return s


def slide_valuation(p, c, chart_ff=None):
    s = _blank(p); v = c["valuation"]
    _kicker_title(s, c["company_name"], v["action_title"])
    _subhead(s, MARGIN, Inches(1.85), Inches(7.2), "Диапазон оценки (football field)")
    _chart(s, chart_ff, MARGIN, Inches(2.35), Inches(7.2), Inches(2.9))
    _txt(s, MARGIN, Inches(5.45), Inches(7.2), Inches(1.4),
         f"Метод: {v.get('method')}. {v.get('justification','')}", size=10.5,
         color=THEME["ink"], spacing=3)
    # правая колонка: апсайд + дисклеймер
    _subhead(s, Inches(8.1), Inches(1.85), Inches(4.6), "Что увеличит оценку")
    _txt(s, Inches(8.1), Inches(2.4), Inches(4.6), Inches(2.0),
         [(x,) for x in v.get("upside_drivers", [])[:5]], size=10.5, color=THEME["ink"], spacing=5)
    _rect(s, Inches(8.1), Inches(4.6), Inches(4.6), Inches(2.3), THEME["hairline"])
    _txt(s, Inches(8.25), Inches(4.7), Inches(4.35), Inches(2.1),
         v.get("disclaimer", ""), size=9, color=THEME["muted"], spacing=2)
    _footer(s, "ФНС / ГИР БО; сравнимые сделки; анализ", 7)
    return s


def slide_highlights(p, c):
    s = _blank(p); hl = c.get("highlights_slide", {})
    _kicker_title(s, c["company_name"], hl.get("action_title", "Точки роста и катализаторы"))
    items = hl.get("catalysts", [])
    n = max(1, len(items))
    cw = (W - 2*MARGIN - Inches(0.3)*(n-1)) / n
    for i, it in enumerate(items[:4]):
        x = MARGIN + i * (cw + Inches(0.3))
        _rect(s, x, Inches(2.1), cw, Inches(4.2), THEME["hairline"])
        _rect(s, x, Inches(2.1), cw, Inches(0.12), THEME["accent"])
        _txt(s, x + Inches(0.2), Inches(2.45), cw - Inches(0.4), Inches(0.7),
             it.get("title", ""), size=14, color=THEME["primary"], bold=True)
        _txt(s, x + Inches(0.2), Inches(3.3), cw - Inches(0.4), Inches(2.8),
             it.get("text", ""), size=11, color=THEME["ink"], spacing=4)
    _footer(s, "анализ", 6)
    return s


def slide_contact(p, c):
    s = _blank(p)
    _rect(s, 0, 0, W, H, THEME["primary"])
    _rect(s, MARGIN, Inches(2.6), Inches(1.4), Pt(3), THEME["accent"])
    _txt(s, MARGIN, Inches(2.9), Inches(11), Inches(1.0), c["contacts"].get("headline", "Следующий шаг"),
         size=30, color=THEME["white"], bold=True)
    _txt(s, MARGIN, Inches(4.1), Inches(11), Inches(1.5), c["contacts"].get("cta", ""),
         size=15, color=THEME["accent"], spacing=6)
    _txt(s, MARGIN, Inches(5.8), Inches(11), Inches(0.6), c["contacts"].get("advisor", ""),
         size=13, color=THEME["white"])
    _txt(s, MARGIN, Inches(6.9), Inches(12), Inches(0.4),
         c["contacts"].get("confidentiality", ""), size=9, color=THEME["muted"])
    return s


def render(content: dict, charts: dict, out_path="teaser.pptx") -> str:
    p = _pres()
    slide_cover(p, content)
    slide_exec_summary(p, content)
    slide_company(p, content)
    slide_market(p, content, chart=charts.get("market"))
    slide_financials(p, content, chart_rev=charts.get("revenue_trend"),
                     chart_margin=charts.get("profit_margin"))
    # катализаторы ДО оценки -> оценка читается как следствие аргументов апсайда
    if content.get("highlights_slide", {}).get("catalysts"):
        slide_highlights(p, content)
    slide_valuation(p, content, chart_ff=charts.get("football_field"))
    slide_contact(p, content)
    p.save(out_path)
    return out_path
